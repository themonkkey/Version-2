"""Extraction + chunking with the two Phase 3 fixes, kept separate from ingest.py
until the pilot proves they help.

Fix 1 (the one the evidence implicates): tables keep their headers.
    ingest.py flattens every table row to " | ".join(cells), so a data row carries no
    record of what its columns mean. A row reading "Paddy | 7275 | 7646" is close to
    meaningless to an embedding model. 60% of the right-document/wrong-page failures
    and 67% of the wrong-document failures are table-heavy chunks.
    Here each data row is serialized against its header row:
        "Paddy productivity | 2024-25: 7275 | 2028-29: 7646"

Fix 2 (correctness, expected to be minor): size caps are actually enforced.
    ingest.py's sub_split only breaks on blank lines, so a single long paragraph is
    never split -- 6,888 chunks exceed their folder cap, and 18.4% of all chunks run
    past Cohere v3's 512-token limit and are silently truncated at embed time.
    Here we fall back to sentence boundaries, then to a hard character split.

Reuses CHUNK_SIZES, OVERLAP, FOLDERS and the boilerplate stripper from ingest.py so
the pilot differs from the baseline only in the two ways above.
"""
import os
import re

import docx
import pdfplumber
from pptx import Presentation

from ingest import CHUNK_SIZES, CORPUS_DIR, OVERLAP, strip_boilerplate

# A table row is only worth emitting if it carries some content beyond its header.
MIN_CELL_CHARS = 2


def _clean(s):
    s = re.sub(r"\s+", " ", (s or "").strip())
    # PDF headers wrap mid-word ("Compl" / "e- tion" on separate rows), so rejoin
    # hyphen-space pairs once the fragments are concatenated: "Comple- tion" -> "Completion"
    return re.sub(r"(\w)-\s+(\w)", r"\1\2", s)


def _header_zone(rows):
    """Index of the first data row.

    These PDFs wrap header text across many sparse rows -- one real table came back as
    17 columns whose header occupied rows 0 through 7 ('Total Land under' / 'Cultivation'
    / 'No. of' / 'Bore-' / 'wells' each on its own row). Treating row 0 as the header
    produces nonsense, so instead find where the data starts: a data row carries several
    populated cells and at least two of them contain digits.
    """
    for i, r in enumerate(rows):
        filled = [c for c in r if c]
        numeric = sum(1 for c in filled if re.search(r"\d", c))
        if len(filled) >= 3 and numeric >= 2:
            return i
    return 1 if len(rows) > 1 else 0


def serialize_table(rows):
    """Turn a table into one line per data row, each bound to its column header.

    Headers are rebuilt per column by concatenating every header-zone fragment in that
    column, which reassembles multi-line headers split across rows by the PDF extractor.
    """
    rows = [[_clean(c) for c in r] for r in rows if r and any(_clean(c) for c in r)]
    if not rows:
        return []
    ncols = max(len(r) for r in rows)
    rows = [r + [""] * (ncols - len(r)) for r in rows]

    hz = _header_zone(rows)
    if hz == 0:
        return [" | ".join(c for c in rows[0] if c)]

    headers = []
    for col in range(ncols):
        frag = [rows[r][col] for r in range(hz) if rows[r][col]]
        headers.append(" ".join(frag).strip())
    # a blank header inherits the nearest populated one to its left, which is how
    # merged/spanning header cells come out of pdfplumber and python-pptx
    last = ""
    for i, h in enumerate(headers):
        if h:
            last = h
        else:
            headers[i] = last

    out = []
    for r in rows[hz:]:
        label = next((c for c in r if len(c) >= MIN_CELL_CHARS and not c.isdigit()), "")
        parts = []
        for i, cell in enumerate(r):
            if not cell or cell == label:
                continue
            head = headers[i] if i < len(headers) else ""
            parts.append(f"{head}: {cell}" if head and head != cell else cell)
        if parts:
            out.append((f"{label} | " if label else "") + " | ".join(parts))
    return out


# ---------------------------------------------------------------- extractors

def extract_pdf_units(path):
    """Prose and tables extracted separately, then recombined.

    The table region is cut out of the page by bounding box before the prose is read,
    so each figure appears exactly once -- in its header-bound serialized form -- rather
    than also surviving as an unlabelled fragment of extract_text() output.
    """
    units = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            try:
                found = page.find_tables()
            except Exception:
                found = []
            prose_page = page
            for tb in found:
                # pdfplumber rejects a crop box that exceeds the page by even a rounding
                # error, and a rejected crop would leave the table's figures in the prose
                # as well as in the serialized rows -- double-counting every number.
                # Clamp to the page bounds so the crop always applies.
                try:
                    px0, ptop, px1, pbot = page.bbox
                    x0, top, x1, bot = tb.bbox
                    box = (max(x0, px0), max(top, ptop), min(x1, px1), min(bot, pbot))
                    if box[2] > box[0] and box[3] > box[1]:
                        prose_page = prose_page.outside_bbox(box)
                except Exception:
                    pass
            text = (prose_page.extract_text() or "") if found else (page.extract_text() or "")
            lines = []
            for tb in found:
                try:
                    lines.extend(serialize_table(tb.extract()))
                except Exception:
                    continue
            if lines:
                text = (text.strip() + "\n\n" + "\n".join(lines)).strip()
            if text.strip():
                units.append({"page": i, "text": text})
    return units


def extract_docx_units(path):
    """One unit per table plus one for the prose, instead of ingest.py's single blob.

    ingest.py concatenates every paragraph and then every table into one giant unit
    with page=None, which detaches tables from the narrative that explains them and
    makes those chunks ineligible for neighbour expansion at retrieval time.
    """
    d = docx.Document(path)
    units = []
    prose = [p.text for p in d.paragraphs if p.text.strip()]
    if prose:
        units.append({"page": None, "text": "\n".join(prose)})
    for t in d.tables:
        rows = [[c.text for c in row.cells] for row in t.rows]
        lines = serialize_table(rows)
        if lines:
            units.append({"page": None, "text": "\n".join(lines)})
    return units


def extract_pptx_units(path):
    prs = Presentation(path)
    units = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs)
                    if t.strip():
                        parts.append(t)
            if shape.has_table:
                rows = [[c.text for c in row.cells] for row in shape.table.rows]
                parts.extend(serialize_table(rows))
        if parts:
            units.append({"page": i, "text": "\n".join(parts)})
    return units


def extract_txt_units(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return [{"page": None, "text": f.read()}]


def extract_units(path):
    ext = path.lower().rsplit(".", 1)[-1]
    try:
        if ext == "pdf":
            return extract_pdf_units(path)
        if ext == "docx":
            return extract_docx_units(path)
        if ext == "pptx":
            return extract_pptx_units(path)
        if ext == "txt":
            return extract_txt_units(path)
    except Exception as e:
        print(f"  ! failed to extract {path}: {e}")
    return []


# ---------------------------------------------------------------- chunking

_SENT = re.compile(r"(?<=[.!?])\s+")


def _hard_split(s, max_chars):
    return [s[i:i + max_chars] for i in range(0, len(s), max_chars)]


def _split_long(block, max_chars):
    """Split an over-long block on sentences, then hard-split anything still over."""
    out, cur = [], ""
    for sent in _SENT.split(block):
        if len(sent) > max_chars:
            if cur:
                out.append(cur.strip())
                cur = ""
            out.extend(_hard_split(sent, max_chars))
            continue
        if len(cur) + len(sent) + 1 > max_chars and cur:
            out.append(cur.strip())
            cur = sent
        else:
            cur = (cur + " " + sent).strip()
    if cur.strip():
        out.append(cur.strip())
    return out


def sub_split(text, page, source, max_chars):
    """Paragraph-first split that actually respects max_chars."""
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        return []
    if len(text) <= max_chars:
        return [{"source": source, "page": page, "text": text}]

    pieces = []
    for para in text.split("\n\n"):
        para = para.strip()
        if not para:
            continue
        pieces.extend(_split_long(para, max_chars) if len(para) > max_chars else [para])

    chunks, cur = [], ""
    for p in pieces:
        if len(cur) + len(p) + 2 > max_chars and cur:
            chunks.append({"source": source, "page": page, "text": cur.strip()})
            # overlap on a word boundary, not mid-word as ingest.py does
            tail = cur[-OVERLAP:]
            tail = tail[tail.find(" ") + 1:] if " " in tail else ""
            cur = (tail + "\n\n" + p).strip()
        else:
            cur = (cur + "\n\n" + p).strip()
    if cur.strip():
        chunks.append({"source": source, "page": page, "text": cur.strip()})
    return chunks


def chunks_for(rel_path):
    """Extract + chunk one corpus-relative file, matching ingest.py's metadata shape.

    No case-study keyword stuffing: ingest.py prepends each tag string four times,
    a TF-IDF-era trick that distorts a dense vector away from the chunk's meaning.
    training_case is the only stratum failing at the document level (docR@10 13-27%),
    so it is a suspect and the pilot drops it.
    """
    folder = rel_path.split("/")[0]
    max_chars = CHUNK_SIZES.get(folder, 1200)
    path = os.path.join(CORPUS_DIR, rel_path)
    out = []
    for u in extract_units(path):
        text = strip_boilerplate(u["text"]) if folder == "case_studies" else u["text"]
        for c in sub_split(text, u["page"], rel_path, max_chars):
            c["folder"] = folder
            out.append(c)
    return out
